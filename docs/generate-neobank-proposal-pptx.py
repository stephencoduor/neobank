#!/usr/bin/env python3
"""Generate NeoBank-Digital-Banking-Proposal.pptx — 16:9 widescreen deck."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "NeoBank-Digital-Banking-Proposal.pptx")

# Savanna colors
PRIMARY = RGBColor(0x2D, 0x6A, 0x4F)
PRIMARY_DARK = RGBColor(0x1B, 0x43, 0x32)
GOLD = RGBColor(0xE9, 0xB9, 0x49)
LIGHT_GREEN = RGBColor(0xD8, 0xF3, 0xDC)
SURFACE = RGBColor(0xFA, 0xFA, 0xF5)
SURFACE_WARM = RGBColor(0xF5, 0xF0, 0xE8)
TEXT_HEADING = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY = RGBColor(0x37, 0x41, 0x51)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_BODY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=16, color=TEXT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0
    return txBox


def add_slide_header(slide, title, subtitle=""):
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), PRIMARY_DARK)
    add_shape(slide, Inches(0), Inches(1.2), SLIDE_W, Pt(4), GOLD)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6),
                 title, 28, WHITE, True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
                     subtitle, 14, LIGHT_GREEN)


def add_footer(slide):
    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(6), Inches(0.4),
                 "Qsoftwares Ltd  |  info@qsoftwares.org  |  +254 710 401 008",
                 10, TEXT_MUTED)
    add_text_box(slide, Inches(10), Inches(7.0), Inches(2.5), Inches(0.4),
                 "CONFIDENTIAL", 10, DANGER, True, PP_ALIGN.RIGHT)


def add_stat_box(slide, left, top, value, label):
    add_shape(slide, left, top, Inches(2.5), Inches(1.4), SURFACE_WARM)
    add_text_box(slide, left, top + Inches(0.15), Inches(2.5), Inches(0.7),
                 str(value), 36, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.85), Inches(2.5), Inches(0.4),
                 label, 13, TEXT_MUTED, False, PP_ALIGN.CENTER)


def add_table(slide, left, top, width, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(1 + len(rows), len(headers), left, top, width, Inches(0.3))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.color.rgb = TEXT_BODY
                para.font.name = "Calibri"
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE

    return table_shape


def add_card(slide, left, top, width, height, title, body, bg_color=WHITE, title_color=PRIMARY):
    shape = add_shape(slide, left, top, width, height, bg_color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                 title, 14, title_color, True)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                 body, 11, TEXT_BODY)


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ──── SLIDE 1: COVER ──────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    # Left green panel
    add_shape(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, PRIMARY_DARK)
    add_shape(slide, Inches(5.5), Inches(0), Pt(5), SLIDE_H, GOLD)
    # Decorative circles
    shape = slide.shapes.add_shape(9, Inches(4.2), Inches(0.5), Inches(1.5), Inches(1.5))  # OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    shape = slide.shapes.add_shape(9, Inches(0.5), Inches(5.5), Inches(1), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x3B, 0x82, 0x5E)
    shape.line.fill.background()

    # Title text on green
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(4.2), Inches(0.5),
                 "Qsoftwares Ltd", 16, GOLD, True)
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(1.2),
                 "Next-Gen\nDigital Banking", 38, WHITE, True)
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(4.2), Inches(0.5),
                 "& Payments Ecosystem", 28, LIGHT_GREEN, True)
    add_text_box(slide, Inches(0.8), Inches(4.3), Inches(4.2), Inches(0.5),
                 "Technical Proposal  |  Gap Analysis  |  Plan", 14, LIGHT_GREEN)
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(4.2), Inches(0.4),
                 "April 2026  |  20-Week Delivery", 14, GOLD, True)

    # Right side — key metrics
    add_text_box(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(0.6),
                 "Platform Readiness", 24, PRIMARY_DARK, True)
    add_stat_box(slide, Inches(6.5), Inches(2.3), "60%", "Backend Ready")
    add_stat_box(slide, Inches(9.3), Inches(2.3), "9", "Payment Providers")
    add_stat_box(slide, Inches(6.5), Inches(4.0), "20", "Weeks to Launch")
    add_stat_box(slide, Inches(9.3), Inches(4.0), "$60K", "Fixed Budget")

    add_text_box(slide, Inches(6.5), Inches(5.8), Inches(6), Inches(0.4),
                 "Mobile-First Financial Operating System", 14, TEXT_MUTED, False, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(6.5), Inches(6.2), Inches(6), Inches(0.4),
                 "for Emerging Markets", 14, TEXT_MUTED, False, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(10), Inches(7.0), Inches(2.5), Inches(0.4),
                 "CONFIDENTIAL", 10, DANGER, True, PP_ALIGN.RIGHT)

    # ──── SLIDE 2: EXECUTIVE SUMMARY ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Executive Summary", "Why We're the Right Partner")
    add_footer(slide)

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.8),
        "We propose building a comprehensive digital banking ecosystem leveraging "
        "our existing Apache Fineract core banking, 9 payment providers, and proven "
        "mobile patterns. ~60% of backend infrastructure is already built.",
        15, TEXT_BODY)

    bullets = [
        "\u2022  Flutter cross-platform app (iOS + Android) with bank-grade security",
        "\u2022  Card issuing via BaaS partner (Marqeta/Stripe) \u2014 virtual + physical",
        "\u2022  Automated KYC/AML via Smile ID (20+ African countries)",
        "\u2022  POS & SoftPOS for merchant payment acceptance",
        "\u2022  P2P payments via phone, alias, and QR codes",
        "\u2022  SOC2/PCI-DSS compliance through architecture + BaaS scope",
    ]
    add_bullets(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(3.0), bullets, 14, TEXT_BODY)

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
        "Key Decision: Flutter + Riverpod + Drift for offline-first, cross-platform mobile.",
        14, PRIMARY, True)

    # ──── SLIDE 3: PLATFORM READINESS ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Platform Readiness", "What We Already Have (Saves 3\u20134 Months)")
    add_footer(slide)

    cards_data = [
        ("Core Banking", "Apache Fineract\nDouble-entry GL\nLoans & Savings\nMulti-currency", PRIMARY_DARK),
        ("9 Payment Providers", "M-Pesa, Airtel, MTN\nFlutterwave, Paystack\nCellulant, AT\nRazorpay, Stripe", PRIMARY),
        ("React Dashboard", "Client CRUD\nLoan/Savings Mgmt\nMobile Money Ops\neTIMS Tax System", PRIMARY_DARK),
        ("Mobile UX Patterns", "50+ Proven Screens\nAuth, Payments, QR\nBiometric Auth\nLoan Applications", PRIMARY),
    ]
    for i, (title, body, tc) in enumerate(cards_data):
        x = Inches(0.8 + i * 3.1)
        add_card(slide, x, Inches(1.6), Inches(2.8), Inches(2.2), title, body, SURFACE, tc)

    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5),
        "Additional: Credit Scoring (PataScore), Multi-Tenant Architecture, eTIMS Tax Compliance",
        13, TEXT_MUTED)

    # Platform diagram — simple text representation
    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
        "Architecture: Flutter App  \u2192  API Gateway  \u2192  Fineract Backend  \u2192  BaaS + Card Issuer + KYC",
        14, PRIMARY, True)

    # ──── SLIDE 4: GAP ANALYSIS ───────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Gap Analysis", "Requirements vs. Current Capabilities")
    add_footer(slide)

    add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3),
        ["Area", "Requirement", "Current State", "Gap"],
        [
            ["Banking", "Tiered Accounts", "Fineract supports client types", "LOW"],
            ["Banking", "KYC/AML Automation", "Document storage only", "HIGH"],
            ["Banking", "Shadow Ledger", "GL exists; pending layer needed", "MEDIUM"],
            ["Cards", "Virtual + Physical Cards", "No card infrastructure", "CRITICAL"],
            ["Cards", "Freeze/Limits/NFC", "None", "CRITICAL"],
            ["Merchant", "POS Terminals + SoftPOS", "No POS infrastructure", "HIGH"],
            ["P2P", "Phone/Alias/QR Payments", "Basic QR + transfers exist", "MEDIUM"],
            ["Security", "SOC2/PCI-DSS", "No compliance framework", "CRITICAL"],
            ["Security", "OAuth2 + JWT", "Basic auth only", "HIGH"],
        ],
        col_widths=[Inches(1.2), Inches(2.8), Inches(4.5), Inches(1.5)])

    # ──── SLIDE 5: CRITICAL SOLUTIONS ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Critical Gap Solutions", "How We Bridge the Gaps")
    add_footer(slide)

    solutions = [
        ("Card Issuing", "BaaS Partner Integration\n\nMarqeta / Stripe Issuing\nVirtual + Physical cards\nTokenization (no raw card data)\nNFC + Chip & PIN via network\n\nEffort: 4\u20135 weeks"),
        ("KYC/AML", "Identity Verification\n\nSmile ID (Africa) / Onfido\nID capture + liveness check\nOCR + AML screening\nAuto-approve or manual queue\n\nEffort: 2\u20133 weeks"),
        ("PCI/SOC2", "Compliance Framework\n\nSAQ-A (outsourced card data)\nTLS 1.3 enforcement\nEncryption at rest (KMS)\nAutomated vuln scanning\n\nEffort: 3\u20134 weeks"),
        ("Merchant/POS", "Payment Acceptance\n\nPAX/Sunmi hardware POS\nMastercard Tap on Phone\nEMVCo QR standard\nInstant settlement engine\n\nEffort: 4\u20135 weeks"),
    ]
    for i, (title, body) in enumerate(solutions):
        x = Inches(0.6 + i * 3.15)
        add_card(slide, x, Inches(1.6), Inches(2.9), Inches(4.5), title, body, SURFACE_WARM, PRIMARY_DARK)

    # ──── SLIDE 6: BAAS STRATEGY ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "BaaS & Sponsor Bank Strategy", "Regional Partners for Banking Rails & Card Issuing")
    add_footer(slide)

    add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3),
        ["Region", "BaaS Partner", "Capabilities", "Card Issuing"],
        [
            ["East Africa", "Cellulant, Stanbic API", "Deposits, payments", "Via Visa/MC partner"],
            ["West Africa", "Flutterwave, Paystack", "Deposits, cards", "Flutterwave virtual"],
            ["Southern Africa", "Stitch, Investec API", "Open banking", "Investec issuing"],
            ["South Asia", "Razorpay X, Setu", "Neo-banking, UPI", "Partner bank cards"],
            ["Southeast Asia", "Brankas, Rapyd", "Wallets", "Rapyd issuing"],
            ["Global", "Marqeta, Stripe Treasury", "Full stack", "Marqeta/Stripe cards"],
        ],
        col_widths=[Inches(2), Inches(3.5), Inches(3.3), Inches(3.5)])

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.0),
        "Architecture Layers:\n"
        "Layer 1: Flutter App + React Dashboard  \u2192  "
        "Layer 2: Kong API Gateway  \u2192  "
        "Layer 3: Fineract Backend  \u2192  "
        "Layer 4: BaaS + Card Issuer + KYC Provider",
        12, TEXT_MUTED)

    # ──── SLIDE 7: TECH STACK ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Technical Stack", "Proven Technologies for Bank-Grade Systems")
    add_footer(slide)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Backend (Extend Existing)", 16, PRIMARY_DARK, True)
    add_table(slide, Inches(0.5), Inches(1.9), Inches(5.8),
        ["Component", "Technology"],
        [
            ["Core Banking", "Apache Fineract (Java 21)"],
            ["API Gateway", "Kong / AWS API Gateway"],
            ["Auth Server", "Keycloak (OAuth2 + JWT)"],
            ["Cache", "Redis"],
            ["Queue", "Apache Kafka"],
            ["Database", "PostgreSQL"],
            ["Monitoring", "Grafana + Prometheus"],
        ],
        col_widths=[Inches(2), Inches(3.8)])

    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Flutter Mobile App", 16, PRIMARY_DARK, True)
    add_table(slide, Inches(6.8), Inches(1.9), Inches(5.8),
        ["Component", "Package"],
        [
            ["Framework", "Flutter 3.x (Dart)"],
            ["State Mgmt", "Riverpod 2.x"],
            ["HTTP", "Dio + Retrofit"],
            ["Local DB", "Drift (SQLite)"],
            ["Biometrics", "local_auth"],
            ["NFC", "nfc_manager"],
            ["QR Codes", "qr_flutter + scanner"],
        ],
        col_widths=[Inches(2), Inches(3.8)])

    # ──── SLIDE 8: IMPLEMENTATION TIMELINE ────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Implementation Timeline", "6 Phases Over 20 Weeks")
    add_footer(slide)

    phases = [
        ("Phase 1", "Foundation\n& Security", "Weeks 1\u20134", "$12K", PRIMARY_DARK),
        ("Phase 2", "Core Banking\nEnhancement", "Weeks 5\u20138", "$12K", PRIMARY),
        ("Phase 3", "Card Issuing\n& Management", "Weeks 9\u201312", "$12K", PRIMARY_DARK),
        ("Phase 4", "Flutter\nMobile App", "Weeks 9\u201316", "$12K", GOLD),
        ("Phase 5", "Merchant\n& POS", "Weeks 13\u201316", "$8K", PRIMARY),
        ("Phase 6", "QA &\nLaunch", "Weeks 17\u201320", "$4K", PRIMARY_DARK),
    ]

    for i, (num, desc, weeks, cost, color) in enumerate(phases):
        x = Inches(0.5 + i * 2.1)
        add_shape(slide, x, Inches(1.8), Inches(1.9), Inches(3.2), color)
        add_text_box(slide, x, Inches(1.9), Inches(1.9), Inches(0.4),
                     num, 14, WHITE, True, PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.1), Inches(2.4), Inches(1.7), Inches(1.0),
                     desc, 14, WHITE, False, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.4), Inches(1.9), Inches(0.4),
                     weeks, 11, LIGHT_GREEN, False, PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.8), Inches(1.9), Inches(0.4),
                     cost, 18, GOLD, True, PP_ALIGN.CENTER)
        # Arrow between phases
        if i < 5:
            add_text_box(slide, x + Inches(1.9), Inches(2.8), Inches(0.3), Inches(0.5),
                         "\u2192", 20, TEXT_MUTED, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.5),
        "Note: Phases 3 & 4 run in parallel (separate teams). Total calendar time: 20 weeks.",
        13, PRIMARY, True)

    # Key deliverables
    add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8),
        "P1: OAuth2, KYC, infra hardening  |  P2: Tiered accounts, shadow ledger, P2P  |  "
        "P3: Card API, management UI  |  P4: Full Flutter app (50+ screens)  |  "
        "P5: POS, SoftPOS, settlement  |  P6: Pen test, SOC2, beta launch",
        11, TEXT_MUTED)

    # ──── SLIDE 9: BUDGET ─────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Budget Breakdown", "$60,000 Fixed Price \u2014 20-Week Delivery")
    add_footer(slide)

    add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3),
        ["Phase", "Duration", "Key Deliverables", "Cost"],
        [
            ["Phase 1: Foundation & Security", "4 weeks", "OAuth2, KYC/AML, infra, BaaS vetting", "$12,000"],
            ["Phase 2: Core Banking Enhancement", "4 weeks", "Tiered accounts, shadow ledger, P2P, FX", "$12,000"],
            ["Phase 3: Card Issuing & Management", "4 weeks", "Card API, management UI, transactions", "$12,000"],
            ["Phase 4: Flutter Mobile App", "8 weeks", "Cross-platform app, 50+ screens", "$12,000"],
            ["Phase 5: Merchant & POS", "4 weeks", "POS, SoftPOS, settlement engine", "$8,000"],
            ["Phase 6: QA, Compliance & Launch", "4 weeks", "Testing, pen test, SOC2, beta", "$4,000"],
            ["TOTAL", "20 weeks", "Complete Digital Banking Platform", "$60,000"],
        ],
        col_widths=[Inches(3.2), Inches(1.5), Inches(5.3), Inches(2.3)])

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
        "Third-party costs (client responsibility): KYC ($200\u2013500/mo), "
        "Card Issuing (per-txn), Cloud ($500\u20132K/mo), POS hardware ($100\u2013300/unit), "
        "SMS ($50\u2013200/mo), Sponsor Bank (negotiated)",
        12, TEXT_MUTED)

    # ──── SLIDE 10: COMPETITIVE ADVANTAGE ─────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_slide_header(slide, "Why Qsoftwares", "Matching Your Brief \u2014 Section 4 Capabilities")
    add_footer(slide)

    add_table(slide, Inches(0.5), Inches(1.6), Inches(12.3),
        ["Your Requirement", "Our Evidence"],
        [
            ["BaaS Strategy & Discovery", "Experience with 9 payment providers across Africa + Asia"],
            ["Financial Backend Engineering", "Apache Fineract ledger + 9 payment webhooks + GL posting"],
            ["Secure Mobile Development", "Biometric auth, payment flows, KYC document capture"],
            ["Infrastructure Hardening", "Cloud deployment, multi-tenant, API security"],
            ["Fintech Project Orchestration", "Multi-app ecosystem (web + 2 mobile) with shared backend"],
            ["Financial QA", "Ledger reconciliation, transaction testing, compliance documentation"],
        ],
        col_widths=[Inches(4), Inches(8.3)])

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
        "Edge Case Matrix maintained for all transaction failure states: "
        "card auth timeout, double STK push, offline P2P, POS battery failure, "
        "FX rate changes, sponsor bank downtime, card fraud.",
        12, TEXT_MUTED)

    # ──── SLIDE 11: NEXT STEPS / CTA ─────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    # Full green panel
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(4.5), PRIMARY_DARK)
    add_shape(slide, Inches(0), Inches(4.5), SLIDE_W, Pt(4), GOLD)

    add_text_box(slide, Inches(2), Inches(0.8), Inches(9), Inches(0.8),
                 "Ready to Build the", 36, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(1.6), Inches(9), Inches(0.8),
                 "Future of Payments?", 36, GOLD, True, PP_ALIGN.CENTER)

    steps = [
        "1.  NDA Signing \u2014 Before disclosing brand name and full specifications",
        "2.  BaaS Discovery Call \u2014 1-week sprint to evaluate 3\u20135 BaaS partners",
        "3.  Technical Deep-Dive \u2014 2-hour session with your CTO",
        "4.  Contract & Kickoff \u2014 Phase 1 begins immediately",
    ]
    add_bullets(slide, Inches(2.5), Inches(2.5), Inches(8), Inches(2.0), steps, 15, LIGHT_GREEN)

    # Contact area
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.6),
                 "Qsoftwares Ltd", 24, PRIMARY, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2), Inches(5.6), Inches(9), Inches(0.4),
                 "info@qsoftwares.org  |  +254 710 401 008  |  www.qsoftwares.org",
                 14, TEXT_BODY, False, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.5),
                 "$60,000 Fixed Price  |  20-Week Delivery  |  Flutter + Fineract + BaaS",
                 16, GOLD, True, PP_ALIGN.CENTER)

    add_text_box(slide, Inches(10), Inches(7.0), Inches(2.5), Inches(0.4),
                 "CONFIDENTIAL", 10, DANGER, True, PP_ALIGN.RIGHT)

    # Save
    prs.save(OUTPUT_PATH)
    size_kb = os.path.getsize(OUTPUT_PATH) / 1024
    print(f"Generated: {OUTPUT_PATH} ({size_kb:.0f} KB)")


if __name__ == "__main__":
    build_pptx()
